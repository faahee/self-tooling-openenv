"""Workflow Agent — Break big goals into steps, execute sequentially, monitor results."""
from __future__ import annotations

import asyncio
import logging
import re
from typing import Callable

from core.llm_core import TaskType

logger = logging.getLogger(__name__)


class WorkflowAgent:
    """Goal-oriented agent that decomposes complex tasks and executes them step by step.

    Example:
        "research competitors, scrape their pricing, save to Excel, email me the report"
        -> Step 1: search for competitors
        -> Step 2: scrape pricing pages
        -> Step 3: save to Excel
        -> Step 4: email report
    """

    # Maximum steps we'll attempt per workflow
    MAX_STEPS = 5

    def __init__(self, config: dict, handle_input_fn: Callable, llm_core, ui_layer, web_agent=None) -> None:
        """
        Args:
            config: JARVIS config dict.
            handle_input_fn: Async function to run a single task (orchestrator.handle_input).
            llm_core: LLMCore for step decomposition.
            ui_layer: UILayer for progress display.
            web_agent: Optional WebAgent for direct search/research.
        """
        self.config = config
        self._handle = handle_input_fn
        self.llm = llm_core
        self.ui = ui_layer
        self.web_agent = web_agent

    # ── Public entry point ────────────────────────────────────────────────────

    async def handle(self, user_input: str, context: dict) -> str:
        """Decompose goal into steps and execute them sequentially.

        Args:
            user_input: High-level goal from user.
            context: Context dict.

        Returns:
            Combined results of all steps.
        """
        self.ui.console.print(
            f"\n  [bold magenta]>> Workflow: decomposing goal into steps...[/bold magenta]"
        )

        steps = await self._decompose(user_input, context)
        if not steps:
            return "Could not decompose the goal into steps. Try being more specific."

        self.ui.console.print(f"  [dim]Goal: {user_input}[/dim]")
        for i, step in enumerate(steps, 1):
            self.ui.console.print(f"  [dim cyan]  Step {i}: {step}[/dim cyan]")

        results: list[dict] = []
        previous_output = ""

        for i, step in enumerate(steps, 1):
            self.ui.console.print(f"\n  [bold cyan]>> Executing step {i}/{len(steps)}: {step}[/bold cyan]")

            # Inject previous step's output into the step if it references "it" or "result"
            enriched_step = self._enrich_step(step, previous_output)

            try:
                # Handle summarize/format steps locally with LLM (they don't need routing)
                result = await asyncio.wait_for(
                    self._execute_step(enriched_step, previous_output), timeout=120
                )
                results.append({"step": i, "command": step, "result": result, "success": True})
                previous_output = result[:6000]
                self.ui.console.print(f"  [green]  Done: {result[:80]}[/green]")
            except asyncio.TimeoutError:
                results.append({"step": i, "command": step, "result": "Timed out", "success": False})
                self.ui.console.print(f"  [red]  Step {i} timed out[/red]")
                # Try to replan remaining steps
                if i < len(steps):
                    replanned = await self._replan(user_input, results, steps[i:])
                    if replanned:
                        steps = steps[:i] + replanned
            except Exception as e:
                results.append({"step": i, "command": step, "result": str(e), "success": False})
                self.ui.console.print(f"  [red]  Step {i} failed: {e}[/red]")

        return self._format_results(user_input, results)

    # ── Step execution ──────────────────────────────────────────────────────

    async def _execute_step(self, step: str, previous_output: str) -> str:
        """Execute a single workflow step, handling LLM-only steps locally."""
        low = step.lower()

        # System monitoring steps — sample CPU/RAM/GPU at intervals and return raw data
        monitor_match = re.match(
            r'monitor system\s+([\w,/ ]+?)\s+every\s+(\d+)\s+seconds?\s+for\s+(\d+)\s+seconds?',
            low
        )
        if monitor_match:
            metrics_raw = monitor_match.group(1).strip()
            interval = int(monitor_match.group(2))
            duration = int(monitor_match.group(3))
            return await self._monitor_system(metrics_raw, interval, duration)

        # Web research/search steps — call web_agent directly to avoid re-routing
        if self.web_agent:
            research_match = re.match(
                r'(?:research|deep search|search)\s+(?:the web )?(?:for|about|on)\s+(.+)',
                low
            )
            if research_match:
                query = research_match.group(1).strip().rstrip("?.")
                if "research" in low:
                    return await self.web_agent.search_and_read(query)
                return await self.web_agent.search(query)

        # Summarize/format/rewrite steps — handle with LLM directly, no routing
        llm_keywords = ("summarize", "rewrite", "format", "create a summary",
                        "write a report", "compile", "draft", "generate report",
                        "create report", "brief report", "professional report",
                        "news report")
        if any(kw in low for kw in llm_keywords) and previous_output:
            prompt = (
                f"{step}\n\n"
                "STRICT RULES:\n"
                "1. ONLY use facts, names, numbers, and events from the 'Source Content' below.\n"
                "2. Do NOT add ANY information from your training data or general knowledge.\n"
                "3. If the source content lacks detail on a topic, omit it entirely.\n"
                "4. Never invent dates, casualty figures, or events not in the source.\n\n"
                "STYLE:\n"
                "- Write as a professional journalist for a major news outlet\n"
                "- Lead with the most important/recent development\n"
                "- Use proper paragraphs, not bullet-point lists of links\n"
                "- Include specific facts, names, and numbers from the source\n"
                "- End with a brief outlook or significance statement\n"
                "- Do NOT list website names or URLs\n\n"
                f"Source Content:\n{previous_output}"
            )
            return await self.llm.generate(prompt, task_type=TaskType.MULTI_STEP_PLAN)

        # ── Format-specific handlers (MUST come BEFORE generic save-to-file) ──

        # Save-as-PDF steps — generate a formatted PDF document
        pdf_match = re.match(
            r'(?:save|export|create|generate|convert|write)\s+.*?'
            r'(?:as pdf|to pdf|pdf)\s*([\w.\-]*\.?pdf)?',
            low
        )
        if not pdf_match:
            pdf_match = re.match(
                r'(?:save|write|export|create|generate)\s+.*?([\w.\-]+\.pdf)',
                low
            )
        if pdf_match and previous_output:
            from pathlib import Path
            filename = pdf_match.group(1) if pdf_match.group(1) else "report.pdf"
            if not filename.endswith(".pdf"):
                filename += ".pdf"
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            self._generate_pdf(str(filepath), previous_output)
            return f"PDF saved to {filepath}"

        # Save-as-PPTX steps — generate a PowerPoint presentation
        pptx_match = re.match(
            r'(?:save|export|create|generate|make)\s+.*?'
            r'(?:as pptx|as presentation|as powerpoint|as slides|to pptx|pptx|presentation|powerpoint|slides)\s*'
            r'([\w.\-]*\.?pptx)?',
            low
        )
        if not pptx_match:
            pptx_match = re.match(
                r'(?:save|create|generate|make)\s+.*?([\w.\-]+\.pptx)',
                low
            )
        if pptx_match and previous_output:
            from pathlib import Path
            filename = pptx_match.group(1) if pptx_match.group(1) else "report.pptx"
            if not filename.endswith(".pptx"):
                filename += ".pptx"
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            self._generate_pptx(str(filepath), previous_output)
            return f"Presentation saved to {filepath}"

        # Save-as-DOCX steps — generate a Word document
        docx_match = re.match(
            r'(?:save|export|create|generate|make)\s+.*?'
            r'(?:as docx|as word|as document|to docx|word doc|docx)\s*'
            r'([\w.\-]*\.?docx)?',
            low
        )
        if not docx_match:
            docx_match = re.match(
                r'(?:save|create|generate|make)\s+.*?([\w.\-]+\.docx)',
                low
            )
        if docx_match and previous_output:
            from pathlib import Path
            filename = docx_match.group(1) if docx_match.group(1) else "report.docx"
            if not filename.endswith(".docx"):
                filename += ".docx"
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            self._generate_docx(str(filepath), previous_output)
            return f"Word document saved to {filepath}"

        # Convert to HTML steps — markdown to styled HTML
        html_match = re.match(
            r'(?:save|export|create|convert|generate)\s+.*?'
            r'(?:as html|to html|html)\s*([\w.\-]*\.?html)?',
            low
        )
        if html_match and previous_output:
            from pathlib import Path
            filename = html_match.group(1) if html_match.group(1) else "report.html"
            if not filename.endswith(".html"):
                filename += ".html"
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            self._generate_html(str(filepath), previous_output)
            return f"HTML saved to {filepath}"

        # Text-to-speech / audio steps — generate MP3 from text
        tts_match = re.match(
            r'(?:convert|save|generate|create|make|read)\s+.*?'
            r'(?:as audio|to audio|as mp3|to mp3|as speech|to speech|audio file|mp3|voice over)\s*'
            r'([\w.\-]*\.?mp3)?',
            low
        )
        if tts_match and previous_output:
            from pathlib import Path
            filename = tts_match.group(1) if tts_match.group(1) else "report.mp3"
            if not filename.endswith(".mp3"):
                filename += ".mp3"
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            self._generate_audio(str(filepath), previous_output)
            return f"Audio saved to {filepath}"

        # Generate chart/graph steps — create a chart image from data
        chart_match = re.match(
            r'(?:create|generate|make|plot|draw|build)\s+.*?'
            r'(?:chart|graph|plot|histogram|bar chart|pie chart|line chart)\s*'
            r'([\w.\-]*\.?(?:png|jpg|svg|pdf))?',
            low
        )
        if chart_match and previous_output:
            from pathlib import Path
            filename = chart_match.group(1) if chart_match.group(1) else "chart.png"
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            self._generate_chart(str(filepath), previous_output, step)
            return f"Chart saved to {filepath}"

        # Generate QR code steps
        qr_match = re.match(
            r'(?:create|generate|make)\s+.*?(?:qr code|qrcode|qr)\s*'
            r'(?:for\s+)?(.+)?',
            low
        )
        if qr_match:
            from pathlib import Path
            qr_data = qr_match.group(1).strip() if qr_match.group(1) else previous_output[:500]
            if qr_data:
                filepath = Path("data") / "qrcode.png"
                filepath.parent.mkdir(parents=True, exist_ok=True)
                self._generate_qr(str(filepath), qr_data)
                return f"QR code saved to {filepath}"

        # Fetch RSS feed steps
        rss_match = re.match(
            r'(?:fetch|get|read|pull|check)\s+.*?(?:rss|feed|atom)\s+(?:from\s+)?(.+)',
            low
        )
        if rss_match:
            url = rss_match.group(1).strip()
            return self._fetch_rss(url)

        # Create calendar event steps
        cal_match = re.match(
            r'(?:create|generate|make|add)\s+.*?'
            r'(?:calendar event|ics|ical|calendar invite|event)\s*'
            r'([\w.\-]*\.?ics)?',
            low
        )
        if cal_match and previous_output:
            from pathlib import Path
            filename = cal_match.group(1) if cal_match.group(1) else "event.ics"
            if not filename.endswith(".ics"):
                filename += ".ics"
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            self._generate_calendar_event(str(filepath), previous_output, step)
            return f"Calendar event saved to {filepath}"

        # OCR — extract text from image
        ocr_match = re.match(
            r'(?:ocr|extract text from|read text from|scan)\s+(.+\.(?:png|jpg|jpeg|bmp|tiff|gif))',
            low
        )
        if ocr_match:
            image_path = ocr_match.group(1).strip()
            return self._ocr_image(image_path)

        # ── Generic save-to-file (LAST — only for plain text formats) ──
        save_match = re.match(
            r'(?:save|write|create file|store)\s+.*?(?:to file|to disk)\s+([\w.\-]+)',
            low
        )
        if not save_match:
            save_match = re.match(
                r'(?:save|write|store)\s+.*?([\w.\-]+\.(?:txt|csv|json|md))',
                low
            )
        if save_match and previous_output:
            from pathlib import Path
            filename = save_match.group(1)
            filepath = Path("data") / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            filepath.write_text(previous_output, encoding="utf-8")
            return f"Saved to {filepath}"

        # Everything else goes through the orchestrator
        return await self._handle(step)

    # ── Step decomposition ────────────────────────────────────────────────────

    async def _decompose(self, goal: str, context: dict | None = None) -> list[str]:
        """Ask LLM to break the goal into concrete executable steps."""
        system_prompt = (context or {}).get("system_prompt", "")
        prompt = f"""Break this goal into concrete steps. Use ONLY these action types:
- "Research the web for <query>" — deep web research (fetches and reads actual articles)
- "Fetch RSS feed from <url>" — pull headlines from an RSS/Atom feed
- "Monitor system <metric> every <N> seconds for <M> seconds" — sample CPU/RAM/GPU/disk at intervals
- "Summarize the research into a professional report" — write a journalist-style summary
- "Save to file <filename>" — write data to disk (txt, csv, json, md)
- "Save as PDF <filename.pdf>" — generate a formatted PDF document
- "Save as PPTX <filename.pptx>" — generate a PowerPoint presentation
- "Save as DOCX <filename.docx>" — generate a Word document
- "Save as HTML <filename.html>" — generate a styled HTML page
- "Convert to audio <filename.mp3>" — generate a spoken MP3 from text
- "Create chart <filename.png>" — generate a bar/pie/line chart from data
- "Create QR code for <data>" — generate a QR code image
- "Create calendar event titled <title>" — generate an .ics event file
- "OCR extract text from <image_file>" — extract text from an image
- "Send email with <subject>" — send email
- "Open <app>" — launch app
- "Analyze <file>" — data analysis

CRITICAL RULES:
- ONLY generate steps the user explicitly asked for. Do NOT add extra output formats.
- If user says "convert to audio", produce ONLY an audio file. Do NOT also create PDF, DOCX, HTML, or PPTX.
- If user says "save as PDF", produce ONLY a PDF. Do NOT also create DOCX, PPTX, or HTML.
- Each workflow should have EXACTLY one output/save step unless user asked for multiple.
- NO setup steps. NO external apps. NO duplicate steps. MAXIMUM {self.MAX_STEPS} steps.
- ALWAYS use "Research the web" (not just "Search") so articles are actually read.
- ALWAYS include a "Summarize" step between research and any output format.

Example — Goal: "find top AI news and save as a report"
1. Research the web for top AI news today
2. Summarize the research into a professional news report
3. Save the summary to file AI_News_Report.txt

Example — Goal: "get the latest news on war and give it in a pdf format"
1. Research the web for latest war news and conflicts today
2. Summarize the research into a professional news report
3. Save as PDF War_News_Report.pdf

Example — Goal: "find latest tech news and convert to audio"
1. Research the web for latest tech news today
2. Summarize the research into a professional news report
3. Convert to audio Latest_Tech_News.mp3

Example — Goal: "research AI and make a presentation"
1. Research the web for latest AI developments
2. Summarize the research into a professional report
3. Save as PPTX AI_Report.pptx

Example — Goal: "monitor CPU and RAM every 5 seconds for 30 seconds and give a summary"
1. Monitor system CPU,RAM every 5 seconds for 30 seconds
2. Summarize the research into a professional report

Goal: {goal}

Output ONLY a numbered list of steps:"""

        response = await self.llm.generate(
            prompt,
            task_type=TaskType.MULTI_STEP_PLAN,
            system_prompt=system_prompt or None,
        )
        return self._parse_steps(response)

    def _parse_steps(self, text: str) -> list[str]:
        """Parse numbered list from LLM output."""
        lines = text.strip().splitlines()
        steps: list[str] = []
        for line in lines:
            line = line.strip()
            # Match "1. step" or "- step" or "* step"
            m = re.match(r'^(?:\d+[.)]\s*|[-*]\s*)(.+)', line)
            if m:
                step = m.group(1).strip()
                if step and len(step) > 3:
                    steps.append(step)
        return steps[:self.MAX_STEPS]

    async def _monitor_system(self, metrics_raw: str, interval: int, duration: int) -> str:
        """Sample system metrics at regular intervals and return raw readings as text.

        Args:
            metrics_raw: Comma-separated metric names e.g. "cpu,ram" or "cpu,ram,gpu".
            interval: Seconds between samples.
            duration: Total monitoring duration in seconds.

        Returns:
            Formatted string of all readings for the summarize step.
        """
        import time as _time
        try:
            import psutil as _psutil
        except ImportError:
            return "psutil not installed — cannot monitor system metrics."

        want_cpu = "cpu" in metrics_raw
        want_ram = "ram" in metrics_raw or "memory" in metrics_raw
        want_gpu = "gpu" in metrics_raw or "vram" in metrics_raw
        want_disk = "disk" in metrics_raw

        readings: list[str] = []
        samples = max(1, duration // interval)

        self.ui.console.print(
            f"  [dim cyan]>> Monitoring {metrics_raw} — {samples} samples "
            f"every {interval}s for {duration}s...[/dim cyan]"
        )

        for i in range(samples):
            parts: list[str] = [f"t={i * interval}s"]
            if want_cpu:
                cpu = _psutil.cpu_percent(interval=None)
                parts.append(f"CPU={cpu}%")
            if want_ram:
                vm = _psutil.virtual_memory()
                parts.append(f"RAM={vm.percent}% ({round(vm.used/1e9,1)}/{round(vm.total/1e9,1)}GB)")
            if want_disk:
                du = _psutil.disk_usage("C:\\")
                parts.append(f"Disk={du.percent}%")
            if want_gpu:
                try:
                    import GPUtil as _gpu
                    gpus = _gpu.getGPUs()
                    if gpus:
                        g = gpus[0]
                        parts.append(f"GPU={g.load*100:.0f}% VRAM={g.memoryUsed:.0f}/{g.memoryTotal:.0f}MB")
                except Exception:
                    pass
            readings.append(" | ".join(parts))
            self.ui.console.print(f"  [dim]  {readings[-1]}[/dim]")
            if i < samples - 1:
                await asyncio.sleep(interval)

        return "System monitoring data:\n" + "\n".join(readings)

    # ── Step enrichment ───────────────────────────────────────────────────────

    def _enrich_step(self, step: str, previous_output: str) -> str:
        """Inject previous result into a step that references it."""
        if not previous_output:
            return step
        low = step.lower()
        ref_words = ("the result", "it", "the output", "that", "the data", "the list")
        if any(w in low for w in ref_words) and previous_output:
            return f"{step} (context from previous step: {previous_output[:200]})"
        return step

    # ── Replanning ────────────────────────────────────────────────────────────

    async def _replan(
        self, goal: str, completed: list[dict], remaining: list[str]
    ) -> list[str]:
        """Ask LLM to adjust remaining steps given what's been done so far."""
        completed_text = "\n".join(
            f"  Step {r['step']}: {r['command']} -> {'OK' if r['success'] else 'FAILED: ' + r['result'][:50]}"
            for r in completed
        )
        prompt = (
            f"Goal: {goal}\n\n"
            f"Completed so far:\n{completed_text}\n\n"
            f"Remaining planned steps (may need adjustment):\n"
            + "\n".join(f"  {i+1}. {s}" for i, s in enumerate(remaining))
            + "\n\nAdjust the remaining steps if needed. Return numbered list only."
        )
        try:
            response = await asyncio.wait_for(
                self.llm.generate(prompt, task_type=TaskType.MULTI_STEP_PLAN), timeout=20
            )
            return self._parse_steps(response)
        except Exception:
            return remaining

    # ── Result formatting ─────────────────────────────────────────────────────

    def _format_results(self, goal: str, results: list[dict]) -> str:
        succeeded = sum(1 for r in results if r["success"])
        failed    = len(results) - succeeded
        lines = [f"Workflow complete: {succeeded}/{len(results)} steps succeeded\n"]

        for r in results:
            icon = "[OK]" if r["success"] else "[FAIL]"
            lines.append(f"{icon} Step {r['step']}: {r['command']}")
            if r["result"] and r["result"].strip():
                # Indent result
                for rline in r["result"].strip().splitlines()[:3]:
                    lines.append(f"       {rline}")
            lines.append("")

        if failed > 0:
            lines.append(f"({failed} step(s) failed — check above for details)")

        return "\n".join(lines)

    # ── Detection ─────────────────────────────────────────────────────────────

    def _generate_pdf(self, filepath: str, content: str) -> None:
        """Generate a formatted PDF from text content using fpdf2."""
        from fpdf import FPDF

        def safe(text: str) -> str:
            """Encode text to latin-1, replacing unsupported chars."""
            return text.encode("latin-1", "replace").decode("latin-1")

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        # Title
        pdf.set_font("Helvetica", "B", 18)
        pdf.cell(0, 12, "JARVIS Report", new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(4)

        # Date
        from datetime import datetime
        pdf.set_font("Helvetica", "I", 10)
        pdf.cell(0, 8, safe(datetime.now().strftime("%B %d, %Y - %I:%M %p")),
                 new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(6)

        # Separator line
        pdf.set_draw_color(100, 100, 100)
        pdf.line(15, pdf.get_y(), 195, pdf.get_y())
        pdf.ln(6)

        # Body content
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                pdf.ln(4)
                continue

            # Detect headings (lines starting with ** or ##)
            if stripped.startswith("**") and stripped.endswith("**"):
                pdf.set_font("Helvetica", "B", 13)
                pdf.multi_cell(0, 7, safe(stripped.strip("*# ")))
                pdf.ln(2)
            elif stripped.startswith("#"):
                pdf.set_font("Helvetica", "B", 14)
                pdf.multi_cell(0, 7, safe(stripped.lstrip("# ")))
                pdf.ln(2)
            elif stripped.startswith(("- ", "* ")):
                pdf.set_font("Helvetica", "", 11)
                pdf.set_x(pdf.l_margin + 8)  # indent
                pdf.multi_cell(pdf.w - pdf.l_margin - pdf.r_margin - 8, 6,
                               safe("- " + stripped.lstrip("-* ")))
            else:
                pdf.set_font("Helvetica", "", 11)
                pdf.multi_cell(0, 6, safe(stripped))

        pdf.output(filepath)

    def _generate_pptx(self, filepath: str, content: str) -> None:
        """Generate a PowerPoint presentation from text content."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from datetime import datetime

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "JARVIS Report"
        slide.placeholders[1].text = datetime.now().strftime("%B %d, %Y")

        # Parse content into sections
        sections: list[dict] = []
        current_title = "Overview"
        current_bullets: list[str] = []

        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            if (stripped.startswith("**") and stripped.endswith("**")) or stripped.startswith("#"):
                if current_bullets:
                    sections.append({"title": current_title, "bullets": current_bullets})
                    current_bullets = []
                current_title = stripped.strip("*# ")
            elif stripped.startswith(("- ", "* ")):
                current_bullets.append(stripped.lstrip("-* ").strip())
            else:
                current_bullets.append(stripped)

        if current_bullets:
            sections.append({"title": current_title, "bullets": current_bullets})

        # Create content slides (max 6 bullets per slide)
        for section in sections:
            for i in range(0, len(section["bullets"]), 6):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = section["title"]
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for j, bullet in enumerate(section["bullets"][i:i + 6]):
                    if j == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                    tf.paragraphs[-1].font.size = Pt(16)

        prs.save(filepath)

    def _generate_docx(self, filepath: str, content: str) -> None:
        """Generate a Word document from text content."""
        from docx import Document
        from docx.shared import Pt
        from datetime import datetime

        doc = Document()
        doc.add_heading("JARVIS Report", level=0)
        doc.add_paragraph(datetime.now().strftime("%B %d, %Y - %I:%M %p")).italic = True
        doc.add_paragraph("")  # spacer

        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            if (stripped.startswith("**") and stripped.endswith("**")) or stripped.startswith("## "):
                doc.add_heading(stripped.strip("*# "), level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped.lstrip("# "), level=1)
            elif stripped.startswith(("- ", "* ")):
                doc.add_paragraph(stripped.lstrip("-* ").strip(), style="List Bullet")
            else:
                doc.add_paragraph(stripped)

        doc.save(filepath)

    def _generate_chart(self, filepath: str, content: str, step: str) -> None:
        """Generate a chart image from text data using matplotlib."""
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        # Try to extract label:value pairs from content
        labels: list[str] = []
        values: list[float] = []

        for line in content.split("\n"):
            stripped = line.strip().lstrip("-* ")
            if not stripped:
                continue
            # Match patterns like "Label: 42", "Label - 42", "Label 42%"
            m = re.match(r'(.+?)[\s:\-–]+\s*([\d,.]+)%?\s*$', stripped)
            if m:
                labels.append(m.group(1).strip()[:30])
                try:
                    values.append(float(m.group(2).replace(",", "")))
                except ValueError:
                    pass

        if not values:
            # Fallback: just create a text-based summary image
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.text(0.5, 0.5, content[:500], transform=ax.transAxes,
                    fontsize=10, verticalalignment='center', horizontalalignment='center',
                    wrap=True)
            ax.axis("off")
            ax.set_title("Data Summary")
        else:
            fig, ax = plt.subplots(figsize=(10, 6))
            # Choose chart type based on step description
            low_step = step.lower()
            if "pie" in low_step:
                ax.pie(values, labels=labels, autopct='%1.1f%%')
                ax.set_title("Distribution")
            elif "line" in low_step:
                ax.plot(labels, values, marker='o', linewidth=2)
                ax.set_title("Trend")
                plt.xticks(rotation=45, ha='right')
            else:
                ax.barh(labels, values, color='steelblue')
                ax.set_title("Comparison")
                ax.set_xlabel("Value")

        plt.tight_layout()
        plt.savefig(filepath, dpi=150, bbox_inches='tight')
        plt.close(fig)

    def _generate_qr(self, filepath: str, data: str) -> None:
        """Generate a QR code image."""
        import qrcode
        qr = qrcode.QRCode(version=1, box_size=10, border=4)
        qr.add_data(data[:2000])
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        img.save(filepath)

    def _generate_html(self, filepath: str, content: str) -> None:
        """Convert markdown content to a styled HTML file."""
        import markdown as md
        from datetime import datetime

        body = md.markdown(content, extensions=["extra", "sane_lists"])
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>JARVIS Report</title>
<style>
  body {{ font-family: 'Segoe UI', Tahoma, sans-serif; max-width: 800px;
         margin: 40px auto; padding: 0 20px; line-height: 1.6; color: #333; }}
  h1 {{ color: #1a1a2e; border-bottom: 2px solid #16213e; padding-bottom: 8px; }}
  h2 {{ color: #16213e; margin-top: 24px; }}
  ul {{ padding-left: 24px; }}
  .meta {{ color: #888; font-style: italic; margin-bottom: 20px; }}
  code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; }}
  pre {{ background: #f4f4f4; padding: 16px; border-radius: 6px; overflow-x: auto; }}
</style>
</head>
<body>
<h1>JARVIS Report</h1>
<p class="meta">Generated on {datetime.now().strftime("%B %d, %Y at %I:%M %p")}</p>
{body}
</body>
</html>"""
        from pathlib import Path
        Path(filepath).write_text(html, encoding="utf-8")

    def _generate_audio(self, filepath: str, content: str) -> None:
        """Generate an MP3 audio file from text using Google TTS."""
        from gtts import gTTS

        # Clean markdown formatting for speech
        clean = re.sub(r'[#*_\-]{2,}', '', content)
        clean = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', clean)  # [text](url) -> text
        clean = clean.strip()

        if len(clean) > 5000:
            clean = clean[:5000]  # gTTS has practical limits

        tts = gTTS(text=clean, lang="en", slow=False)
        tts.save(filepath)

    def _fetch_rss(self, url: str) -> str:
        """Fetch and parse an RSS/Atom feed, return formatted text."""
        import feedparser

        if not url.startswith(("http://", "https://")):
            url = "https://" + url

        feed = feedparser.parse(url)

        if feed.bozo and not feed.entries:
            return f"Could not parse RSS feed: {url}"

        title = feed.feed.get("title", "RSS Feed")
        lines = [f"# {title}\n"]

        for entry in feed.entries[:15]:
            entry_title = entry.get("title", "No title")
            link = entry.get("link", "")
            published = entry.get("published", "")
            summary = entry.get("summary", "")
            # Strip HTML from summary
            summary = re.sub(r'<[^>]+>', '', summary)[:200]

            lines.append(f"**{entry_title}**")
            if published:
                lines.append(f"  Published: {published}")
            if summary:
                lines.append(f"  {summary}")
            if link:
                lines.append(f"  Link: {link}")
            lines.append("")

        return "\n".join(lines)

    def _generate_calendar_event(self, filepath: str, content: str, step: str) -> None:
        """Generate an .ics calendar event file."""
        from icalendar import Calendar, Event
        from datetime import datetime, timedelta
        import uuid

        cal = Calendar()
        cal.add("prodid", "-//JARVIS//EN")
        cal.add("version", "2.0")

        event = Event()

        # Try to extract event details from step/content
        # Default to a generic event based on the content
        title_match = re.search(r'(?:titled?|called?|named?|about)\s+"?([^"]+)"?', step, re.I)
        event_title = title_match.group(1).strip() if title_match else "JARVIS Event"

        event.add("summary", event_title)
        event.add("description", content[:1000])
        event.add("dtstart", datetime.now() + timedelta(hours=1))
        event.add("dtend", datetime.now() + timedelta(hours=2))
        event.add("uid", str(uuid.uuid4()))

        cal.add_component(event)

        from pathlib import Path
        Path(filepath).write_bytes(cal.to_ical())

    def _ocr_image(self, image_path: str) -> str:
        """Extract text from an image using Tesseract OCR."""
        from pathlib import Path

        path = Path(image_path)
        if not path.exists():
            return f"Image not found: {image_path}"

        try:
            import pytesseract
            from PIL import Image
            img = Image.open(str(path))
            text = pytesseract.image_to_string(img)
            return text.strip() if text.strip() else "No text detected in image."
        except Exception as e:
            return f"OCR failed: {e}"

    # ── Detection ─────────────────────────────────────────────────────────────

    @staticmethod
    def is_workflow(user_input: str) -> bool:
        """Heuristic: detect if input describes a multi-step workflow."""
        low = user_input.lower()
        non_workflows = (
            r"\b(?:run|execute)\s+powershell\b",
            r"\bpowershell command\b",
            r"\b(?:run|execute)\s+cmd\b",
            r"\b(?:run|execute)\s+command\b",
        )
        if any(re.search(pattern, low) for pattern in non_workflows):
            return False
        simple_file_ops = (
            r"\b(?:create|make)\s+(?:a\s+)?folder\b",
            r"\b(?:create|make)\s+(?:a\s+)?directory\b",
            r"\b(?:move|rename|copy|delete)\s+(?:a\s+)?(?:file|folder|directory)\b",
        )
        if any(re.search(pattern, low) for pattern in simple_file_ops):
            return False
        # Multiple actions chained with then/then/and/after
        chain_words = ("then", "after that", "next", "finally", "also", "and then", "afterwards")
        has_chain = sum(1 for w in chain_words if w in low) >= 1

        # Multiple commas or semicolons suggesting a list of tasks
        has_list = user_input.count(",") >= 2 or user_input.count(";") >= 1

        # Step-like goal phrases
        goal_phrases = (
            "step by step", "one by one", "in order",
            "research and", "collect and", "analyze and",
            "scrape and save", "find and email", "check and notify",
        )
        has_goal_phrase = any(p in low for p in goal_phrases)

        # Output-format patterns: user wants info + a specific deliverable
        # e.g. "get news on war and give it in a pdf", "find competitors and save to excel"
        _INFO_VERBS = ("news", "find", "get", "search", "research", "look up",
                       "fetch", "collect", "tell me", "i want", "give me",
                       "show me", "latest", "recent")
        _OUTPUT_FORMATS = ("pdf", "excel", "csv", "report", "document",
                          "word doc", "docx", "spreadsheet", "pptx", "powerpoint",
                          "presentation", "slides", "html", "web page",
                          "mp3", "audio", "voice over", "chart", "graph", "plot",
                          "qr code", "qrcode", "calendar event", "ics",
                          "email it", "email me",
                          "save it", "save to", "send it", "send me",
                          "as a file", "to a file", "in a file",
                          "create a report", "write a report", "make a report")
        has_info = any(v in low for v in _INFO_VERBS)
        has_output = any(f in low for f in _OUTPUT_FORMATS)
        has_info_plus_output = has_info and has_output

        # Conjunction pattern: "<task> and <deliver>" where deliver implies output
        _DELIVER_VERBS = ("and give", "and save", "and send", "and email",
                          "and create", "and export", "and convert",
                          "and make", "and write", "and store")
        has_deliver = any(d in low for d in _DELIVER_VERBS)

        # Multi-app or multi-action OS command: "open X, Y and do Z"
        # One comma + "and" + a non-app action verb = multi-step
        _ACTION_VERBS = ("take a screenshot", "screenshot", "send", "email",
                         "save", "open", "launch", "search", "play", "mute",
                         "close", "kill", "type", "write", "create")
        has_multi_action = (
            user_input.count(",") >= 1
            and sum(1 for v in _ACTION_VERBS if v in low) >= 2
        )

        return (has_chain and has_list) or has_goal_phrase or has_info_plus_output or has_deliver or has_multi_action
